from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List

# Optional dependencies:
# - python-docx
# - openpyxl
# - python-pptx
# - pypdf


def extract_text_and_outline(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            return _extract_docx(path, limits)
        if ext == ".xlsx":
            return _extract_xlsx(path, limits)
        if ext == ".pptx":
            return _extract_pptx(path, limits)
        if ext == ".pdf":
            return _extract_pdf(path, limits)
        if ext in {".md", ".txt", ".sql"}:
            return _extract_text(path, limits)
        # unknown
        return {"headings": [], "preview": ""}
    except Exception:
        return {"headings": [], "preview": ""}


def _extract_text(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    txt = path.read_text(encoding="utf-8", errors="ignore")
    preview = txt[: limits["max_extract_chars"]]
    headings = []
    for line in preview.splitlines():
        if line.strip().startswith("#"):
            headings.append(line.strip().lstrip("#").strip())
        if len(headings) >= limits["max_headings"]:
            break
    return {"headings": headings, "preview": preview}


def _extract_docx(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    from docx import Document  # type: ignore

    doc = Document(str(path))
    headings: List[str] = []
    preview_parts: List[str] = []
    para_count = 0

    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue

        style = (p.style.name or "") if p.style else ""
        if style.lower().startswith("heading") or "見出し" in style:
            headings.append(text)
        else:
            if para_count < limits["max_preview_paragraphs"]:
                preview_parts.append(text)
                para_count += 1

        if len(headings) >= limits["max_headings"]:
            # still collect preview but bounded
            pass
        if sum(len(x) for x in preview_parts) >= limits["max_extract_chars"]:
            break

    preview = "\n".join(preview_parts)[: limits["max_extract_chars"]]
    return {"headings": headings[: limits["max_headings"]], "preview": preview}


def _extract_xlsx(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    headings = [f"Sheet: {name}" for name in wb.sheetnames][: limits["max_headings"]]

    preview_parts: List[str] = []
    cells_collected = 0

    # preview: top-left cells from first few sheets
    for sidx, name in enumerate(wb.sheetnames[:5]):
        ws = wb[name]
        preview_parts.append(f"[{name}]")
        for row in ws.iter_rows(min_row=1, max_row=30, min_col=1, max_col=12, values_only=True):
            row_vals = []
            for v in row:
                if v is None:
                    continue
                sv = str(v).strip()
                if not sv:
                    continue
                row_vals.append(sv)
                cells_collected += 1
                if cells_collected >= limits["max_preview_cells"]:
                    break
            if row_vals:
                preview_parts.append(" | ".join(row_vals))
            if cells_collected >= limits["max_preview_cells"]:
                break
        if cells_collected >= limits["max_preview_cells"]:
            break

    preview = "\n".join(preview_parts)[: limits["max_extract_chars"]]
    return {"headings": headings, "preview": preview}


def _extract_pptx(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    headings: List[str] = []
    preview_parts: List[str] = []

    for idx, slide in enumerate(prs.slides[: limits["max_preview_slides"]]):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        if title:
            headings.append(f"Slide {idx+1}: {title}")
            preview_parts.append(f"[Slide {idx+1}] {title}")
        else:
            preview_parts.append(f"[Slide {idx+1}]")

        # bullets / text frames
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            txt = (shape.text_frame.text or "").strip()
            if txt:
                preview_parts.append(txt)

        if sum(len(x) for x in preview_parts) >= limits["max_extract_chars"]:
            break

    preview = "\n".join(preview_parts)[: limits["max_extract_chars"]]
    return {"headings": headings[: limits["max_headings"]], "preview": preview}


def _extract_pdf(path: Path, limits: Dict[str, int]) -> Dict[str, Any]:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    preview_parts: List[str] = []
    headings: List[str] = []

    for i, page in enumerate(reader.pages[:10]):  # PoC: first 10 pages
        txt = (page.extract_text() or "").strip()
        if not txt:
            continue
        # naive: first line as a heading candidate
        first_line = txt.splitlines()[0].strip() if txt.splitlines() else ""
        if first_line:
            headings.append(f"Page {i+1}: {first_line[:120]}")
        preview_parts.append(txt)

        if sum(len(x) for x in preview_parts) >= limits["max_extract_chars"]:
            break

    preview = "\n".join(preview_parts)[: limits["max_extract_chars"]]
    return {"headings": headings[: limits["max_headings"]], "preview": preview}
